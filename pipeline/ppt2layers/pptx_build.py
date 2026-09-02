"""python-pptx 组装辅助：精确文本框（无自动换行/无内边距/固定行距）、图片层、PowerPoint 导出 PDF、PDF 渲染。"""
import os, subprocess, numpy as np
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from .common import W0, H0, ROOT

SLIDE_W_EMU, SLIDE_H_EMU = 12192000, 6858000
SLIDE_W_PT = SLIDE_W_EMU / 12700.0          # 960pt
PX2PT = SLIDE_W_PT / W0                      # 1px → pt
PT2PX = 1 / PX2PT
def px2emu(px): return int(round(px * SLIDE_W_EMU / W0))
def pt2emu(pt): return int(round(pt * 12700))

NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
         "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
         "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}

def new_prs():
    prs = Presentation(); prs.slide_width = Emu(SLIDE_W_EMU); prs.slide_height = Emu(SLIDE_H_EMU)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_picture(slide, path, box_px, name=None):
    x0, y0, x1, y1 = box_px
    pic = slide.shapes.add_picture(path, Emu(px2emu(x0)), Emu(px2emu(y0)), Emu(px2emu(x1 - x0)), Emu(px2emu(y1 - y0)))
    if name: pic.name = name
    return pic

def _hex(c): return c.lstrip("#").upper()

def add_textbox(slide, lines, x_px, y_px, w_px, size_pt, font="微软雅黑", bold=False, color="#000000",
                align="left", spc_pt=0.0, pitch_pt=None, name=None, glow=None, h_px=None, latin_font=None, outline=None, indents_px=None):
    """lines: 每行一个段落。x_px/y_px 为文本框左上角（像素，图像坐标）。pitch_pt: 固定行距(磅)，None=单倍。
    glow: None 或 {"radius_pt":..., "color":"#RRGGBB", "alpha":0~1}"""
    if isinstance(pitch_pt, (list, tuple)): pitch_pt = list(pitch_pt) or None
    if h_px is None:
        if isinstance(pitch_pt, list):
            h_px = (sum(pitch_pt) + pitch_pt[0]) * PT2PX + 4
        else:
            h_px = (pitch_pt if pitch_pt else size_pt * 1.2) * PT2PX * max(1, len(lines)) + 4
    tb = slide.shapes.add_textbox(Emu(px2emu(x_px)), Emu(px2emu(y_px)), Emu(px2emu(w_px)), Emu(px2emu(h_px)))
    if name: tb.name = name
    txBody = tb.text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    for k in list(bodyPr.attrib): del bodyPr.attrib[k]
    bodyPr.set("wrap", "none"); bodyPr.set("lIns", "0"); bodyPr.set("tIns", "0"); bodyPr.set("rIns", "0"); bodyPr.set("bIns", "0")
    bodyPr.set("anchor", "t"); bodyPr.set("rtlCol", "0")
    for ch in list(bodyPr): bodyPr.remove(ch)
    etree.SubElement(bodyPr, qn("a:noAutofit"))
    # 清掉默认段落
    for p in txBody.findall(qn("a:p")): txBody.remove(p)
    algn = {"left": "l", "center": "ctr", "right": "r"}[align]
    for li, line in enumerate(lines):
        runs = [{"t": line}] if isinstance(line, str) else line["runs"]
        p = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr")); pPr.set("algn", algn)
        if indents_px and li < len(indents_px) and indents_px[li]:
            pPr.set("marL", "0"); pPr.set("indent", str(px2emu(indents_px[li])))
        if pitch_pt:
            # 逐段行距：段 li 的 lnSpc 决定它与上一行的间隙（列表=各间隙实测值；首段用首间隙，只影响框内起点，y_top 会吸收）
            pv = (pitch_pt[0] if li == 0 else pitch_pt[min(li - 1, len(pitch_pt) - 1)]) if isinstance(pitch_pt, list) else pitch_pt
            ln = etree.SubElement(pPr, qn("a:lnSpc")); sp = etree.SubElement(ln, qn("a:spcPts")); sp.set("val", str(int(round(pv * 100))))
        sb = etree.SubElement(pPr, qn("a:spcBef")); e = etree.SubElement(sb, qn("a:spcPts")); e.set("val", "0")
        sa = etree.SubElement(pPr, qn("a:spcAft")); e = etree.SubElement(sa, qn("a:spcPts")); e.set("val", "0")
        for run in runs:
          r = etree.SubElement(p, qn("a:r"))
          rPr = etree.SubElement(r, qn("a:rPr")); rPr.set("lang", "zh-CN"); rPr.set("altLang", "en-US")
          rb = bold if run.get("bold") is None else run["bold"]
          rcolor = run.get("color") or color
          rPr.set("sz", str(int(round(size_pt * 100)))); rPr.set("b", "1" if rb else "0"); rPr.set("dirty", "0")
          if abs(spc_pt) > 1e-3: rPr.set("spc", str(int(round(spc_pt * 100))))
          if outline and outline.get("mode", "ln") == "ln":
            ln = etree.SubElement(rPr, qn("a:ln")); ln.set("w", str(pt2emu(outline["width_pt"])))
            lf = etree.SubElement(ln, qn("a:solidFill")); lc = etree.SubElement(lf, qn("a:srgbClr")); lc.set("val", _hex(outline["color"]))
            if outline.get("alpha", 1.0) < 0.999:
                la = etree.SubElement(lc, qn("a:alpha")); la.set("val", str(int(round(outline["alpha"] * 100000))))
            rnd = etree.SubElement(ln, qn("a:round"))
          sf = etree.SubElement(rPr, qn("a:solidFill")); c = etree.SubElement(sf, qn("a:srgbClr")); c.set("val", _hex(rcolor))
          if glow or (outline and outline.get("mode") == "shadow"):
            eff = etree.SubElement(rPr, qn("a:effectLst"))
            if outline and outline.get("mode") == "shadow":   # 软阴影模式
                sh = etree.SubElement(eff, qn("a:outerShdw")); sh.set("blurRad", str(pt2emu(outline["blur_pt"]))); sh.set("dist", "0"); sh.set("dir", "0"); sh.set("algn", "ctr"); sh.set("rotWithShape", "0")
                sc = etree.SubElement(sh, qn("a:srgbClr")); sc.set("val", _hex(outline["color"]))
                al = etree.SubElement(sc, qn("a:alpha")); al.set("val", str(int(round(outline.get("alpha", 0.8) * 100000))))
            if glow:
                g = etree.SubElement(eff, qn("a:glow")); g.set("rad", str(pt2emu(glow["radius_pt"])))
                gc = etree.SubElement(g, qn("a:srgbClr")); gc.set("val", _hex(glow.get("color", "#FFFFFF")))
                al = etree.SubElement(gc, qn("a:alpha")); al.set("val", str(int(round(glow.get("alpha", 0.4) * 100000))))
          lf = latin_font or font
          for tag, tf in (("a:latin", lf), ("a:ea", font), ("a:cs", font)):
            el = etree.SubElement(rPr, qn(tag)); el.set("typeface", tf)
          t = etree.SubElement(r, qn("a:t")); t.text = run["t"]
        # 段落结束属性：保持字号一致（避免空段落默认字号影响行高）
        ep = etree.SubElement(p, qn("a:endParaRPr")); ep.set("lang", "zh-CN"); ep.set("sz", str(int(round(size_pt * 100))))
    return tb

import time as _time

def _pp_lock(acquire=True, lock_dir="/tmp/ppt2layers_powerpoint.lock"):
    if not acquire:
        try: os.rmdir(lock_dir)
        except Exception: pass
        return True
    for _ in range(360):           # 最多等 30 分钟
        try:
            os.mkdir(lock_dir); return True
        except FileExistsError:
            # 陈旧锁（超过 30 分钟）直接抢占
            try:
                if _time.time() - os.path.getmtime(lock_dir) > 1800:
                    os.rmdir(lock_dir); continue
            except Exception: pass
            _time.sleep(5)
    return False

RENDERER = os.environ.get("SLIDELIFT_RENDERER", "powerpoint")   # powerpoint | libreoffice

# LibreOffice 模式下"标定稿"用的免费替身字体（服务器无微软/华文字体授权）。
# 成品 assemble 仍写原字体名——终端用户的 Windows/Office 上有真字体；
# 中文字形在几乎所有字体中都是整 em 等宽，替身标定出的字号/位置可迁移，残差主要在拉丁字符。
CAL_FONT_MAP = {
    "微软雅黑": "Source Han Sans CN", "黑体": "Source Han Sans CN",
    "楷体": "LXGW WenKai", "华文行楷": "LXGW WenKai",
    "宋体": "Source Han Serif CN",
    "Arial": "Liberation Sans",     # LibreOffice 自带，与 Arial 度量兼容
}
def cal_font(name):
    """标定/参考渲染场景的字体名：LO 模式映射到本机可见的替身，PowerPoint 模式原样。"""
    return CAL_FONT_MAP.get(name, name) if RENDERER == "libreoffice" else name

def _soffice_bin():
    cand = os.environ.get("SLIDELIFT_SOFFICE")
    if cand and os.path.exists(cand): return cand
    for p in ("/Applications/LibreOffice.app/Contents/MacOS/soffice", "/usr/bin/soffice", "/usr/local/bin/soffice"):
        if os.path.exists(p): return p
    import shutil as _sh
    found = _sh.which("soffice")
    if found: return found
    raise RuntimeError("找不到 LibreOffice（soffice）——设 SLIDELIFT_SOFFICE 环境变量")

def export_pdf_lo(pptx_path, pdf_path, timeout=300):
    """LibreOffice 无头渲染（服务器路径）。独立 UserInstallation 目录避免实例互踩；
    soffice 输出文件名固定为 <stem>.pdf，按需改名到目标路径。"""
    pptx_path = os.path.abspath(pptx_path); pdf_path = os.path.abspath(pdf_path)
    if not os.path.exists(pptx_path): raise FileNotFoundError(pptx_path)
    if os.path.exists(pdf_path): os.remove(pdf_path)
    outdir = os.path.dirname(pdf_path)
    profile = f"/tmp/slidelift_lo_{os.getpid()}"
    r = subprocess.run([_soffice_bin(), "--headless", "--norestore", "--nolockcheck",
                        f"-env:UserInstallation=file://{profile}",
                        "--convert-to", "pdf", "--outdir", outdir, pptx_path],
                       capture_output=True, text=True, timeout=timeout)
    produced = os.path.join(outdir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if not os.path.exists(produced):
        raise RuntimeError(f"LibreOffice export failed: rc={r.returncode} {r.stderr[:300]}")
    if produced != pdf_path: os.replace(produced, pdf_path)
    return pdf_path

def export_pdf(pptx_path, pdf_path, timeout=900):
    """pptx → PDF。默认本机 Microsoft PowerPoint（AppleScript）；
    SLIDELIFT_RENDERER=libreoffice 时走 LibreOffice 无头（服务器部署路径）。"""
    if RENDERER == "libreoffice":
        return export_pdf_lo(pptx_path, pdf_path)
    pptx_path = os.path.abspath(pptx_path); pdf_path = os.path.abspath(pdf_path)
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(pptx_path)   # 千万别让 PowerPoint 去打开不存在的文件：会弹模态框卡死后续所有 AppleEvent
    if os.path.exists(pdf_path): os.remove(pdf_path)
    script = f'''
with timeout of {timeout} seconds
tell application "Microsoft PowerPoint"
	open POSIX file "{pptx_path}"
	set pres to active presentation
	save pres in POSIX file "{pdf_path}" as save as PDF
	close pres saving no
end tell
end timeout'''
    if not _pp_lock(True):
        raise RuntimeError("PowerPoint busy: lock timeout")
    try:
        last = None
        for attempt in range(2):
            r = subprocess.run(["perl", "-e", f"alarm {timeout + 60}; exec @ARGV", "osascript", "-e", script], capture_output=True, text=True)
            if r.returncode == 0 and os.path.exists(pdf_path):
                return pdf_path
            last = f"rc={r.returncode} {r.stderr[:400]}"
            _time.sleep(8)
        raise RuntimeError(f"PowerPoint export failed: {last}")
    finally:
        _pp_lock(False)

def render_pdf(pdf_path, width=W0):
    import pymupdf
    doc = pymupdf.open(pdf_path); out = []
    for pg in doc:
        z = width / pg.rect.width
        pix = pg.get_pixmap(matrix=pymupdf.Matrix(z, z), alpha=False, colorspace=pymupdf.csRGB)
        arr = np.frombuffer(pix.samples, np.uint8).reshape(pix.height, pix.width, 3).copy()
        out.append(arr)
    return out
